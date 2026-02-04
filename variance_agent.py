import pandas as pd

# Step 1: Read file as raw text (single column)
df_raw = pd.read_csv("fpa_data.csv", header=None)

# Step 2: Split the column by comma
df = df_raw[0].str.split(",", expand=True)

# Step 3: Drop the header row (row 0)
df = df.iloc[1:]

# Step 4: Assign correct column names
df.columns = ["Month", "Entity", "Account", "Type", "Amount"]

# Step 5: Convert Amount to number
df["Amount"] = pd.to_numeric(df["Amount"])

print("COLUMNS:", df.columns.tolist())
print(df.head())

# Step 6: Split Actual vs Budget
actual = df[df["Type"] == "Actual"]
budget = df[df["Type"] == "Budget"]

# Step 7: Merge Actual and Budget
merged = actual.merge(
    budget,
    on=["Month", "Entity", "Account"],
    suffixes=("_Actual", "_Budget")
)

# Step 8: Calculate variance
merged["Variance"] = merged["Amount_Actual"] - merged["Amount_Budget"]

# Step 9: Sort worst first
ranked = merged.sort_values("Variance")

print("\nFINAL VARIANCE TABLE:")
print(ranked)




from openai import OpenAI
import textwrap

client = OpenAI()

# Convert variance table to text for AI
variance_text = ranked.to_string(index=False)

prompt = f"""
You are an FP&A Controller writing management commentary.

Tasks:
1. Identify the top unfavorable variances
2. Explain drivers using finance logic (volume, price, mix, cost)
3. Use concise CFO-ready language
4. Do NOT invent assumptions
5. Use numbers

Variance Table:
{variance_text}
"""

commentary = """
January results are below budget primarily driven by APAC revenue underperformance
(-80k vs budget) and EMEA revenue shortfall (-50k), indicating lower-than-expected
volumes. AMER COGS exceeded budget by 20k, suggesting cost inflation pressures,
partially offset by revenue overperformance. SG&A remains broadly controlled.
"""

print("\nAI COMMENTARY:\n")
print(textwrap.fill(commentary, 100))




from pptx import Presentation

# Create presentation
prs = Presentation()

# -----------------------------
# Slide 1: Title
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "January 2025 FP&A Variance Analysis"
slide.placeholders[1].text = "Actual vs Budget"

# -----------------------------
# Slide 2: Executive Summary
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Executive Summary"
slide.placeholders[1].text = commentary.strip()

# -----------------------------
# Slide 3: Key Variances
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Key Variances (Actual vs Budget)"

bullets = []
for _, row in ranked.head(5).iterrows():
    bullets.append(
        f"{row['Entity']} {row['Account']}: {int(row['Variance']):,}"
    )

slide.placeholders[1].text = "\n".join(bullets)

# Save presentation
prs.save("FP&A_Variance_Report_Jan_2025.pptx")

