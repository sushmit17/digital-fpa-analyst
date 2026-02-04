import streamlit as st
import pandas as pd
from pptx import Presentation


# -----------------------------
# FP&A ENGINE (REUSED LOGIC)
# -----------------------------

def run_variance_engine(uploaded_file):

    # Read file as raw text
    df_raw = pd.read_csv(uploaded_file, header=None)

    # Split columns
    df = df_raw[0].str.split(",", expand=True)

    # Drop header
    df = df.iloc[1:]

    # Assign names
    df.columns = ["Month", "Entity", "Account", "Type", "Amount"]

    # Convert Amount
    df["Amount"] = pd.to_numeric(df["Amount"])

    # Split
    actual = df[df["Type"] == "Actual"]
    budget = df[df["Type"] == "Budget"]

    # Merge
    merged = actual.merge(
        budget,
        on=["Month", "Entity", "Account"],
        suffixes=("_Actual", "_Budget")
    )

    # Variance
    merged["Variance"] = merged["Amount_Actual"] - merged["Amount_Budget"]

    ranked = merged.sort_values("Variance")

    return ranked


# -----------------------------
# PPT GENERATOR
# -----------------------------

def generate_ppt(ranked, commentary):

    prs = Presentation()

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "FP&A Variance Report"
    slide.placeholders[1].text = "Automated Analysis"

    # Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    slide.placeholders[1].text = commentary

    # Variances
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Top Variances"

    bullets = []
    for _, row in ranked.head(5).iterrows():
        bullets.append(
            f"{row['Entity']} {row['Account']}: {int(row['Variance']):,}"
        )

    slide.placeholders[1].text = "\n".join(bullets)

    file_name = "FP&A_Report.pptx"
    prs.save(file_name)

    return file_name


# -----------------------------
# STREAMLIT UI
# -----------------------------

st.set_page_config(page_title="Digital FP&A Analyst")

st.title("📊 Digital FP&A Analyst")
st.write("Upload your Actual vs Budget file and generate a management report.")


uploaded_file = st.file_uploader(
    "Upload CSV file",
    type=["csv"]
)


if uploaded_file:

    st.success("File uploaded successfully.")

    if st.button("Run Variance Analysis"):

        with st.spinner("Running FP&A engine..."):

            ranked = run_variance_engine(uploaded_file)

            st.subheader("Variance Preview")
            st.dataframe(ranked.head(10))

            # Mocked AI commentary
            commentary = """
January results are below budget primarily driven by APAC and EMEA revenue shortfalls.
Cost pressures were observed in AMER COGS. SG&A remains controlled.
"""

            ppt_file = generate_ppt(ranked, commentary)

            st.success("Report generated!")

            with open(ppt_file, "rb") as f:
                st.download_button(
                    label="Download PowerPoint",
                    data=f,
                    file_name=ppt_file,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
